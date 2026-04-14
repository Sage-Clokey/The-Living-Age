"""
Build PowerPoint progress report slides for BME 129C Capstone.
Uses python-pptx to create image-heavy presentation.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FIGURES = Path(__file__).resolve().parent / "figures"
OUT = Path(__file__).resolve().parent / "progress_report_slides.pptx"

# Colors
GREEN_DARK = RGBColor(0x2d, 0x6a, 0x4f)
GREEN_MID = RGBColor(0x52, 0xb7, 0x88)
GREEN_LIGHT = RGBColor(0x95, 0xd5, 0xb2)
GOLD = RGBColor(0xe9, 0xc4, 0x6a)
RED = RGBColor(0xe6, 0x39, 0x46)
WHITE = RGBColor(0xff, 0xff, 0xff)
OFF_WHITE = RGBColor(0xf5, 0xf5, 0xf5)
DARK_BG = RGBColor(0x0d, 0x11, 0x17)
PANEL_BG = RGBColor(0x16, 0x1b, 0x22)
GRAY = RGBColor(0x8b, 0x94, 0x9e)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(slide, bullets, left, top, width, height,
                              font_size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, sub_bullets) in enumerate(bullets):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0

        for sub in sub_bullets:
            p2 = tf.add_paragraph()
            p2.text = sub
            p2.font.size = Pt(font_size - 2)
            p2.font.color.rgb = GRAY
            p2.font.name = "Calibri"
            p2.space_after = Pt(4)
            p2.level = 1

    return txBox


def add_green_bar(slide, top=Inches(0), height=Inches(0.06)):
    """Add a thin green accent bar across the top."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, Inches(13.33), height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN_MID
    shape.line.fill.background()


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # =========================================================================
    # SLIDE 1: Title
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                 "Living Systems as Decentralized Economies",
                 font_size=40, color=GREEN_MID, bold=True)

    add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
                 "Why the Molecular Biologist Should Be an Entrepreneur, Not a Central Planner",
                 font_size=22, color=GOLD)

    add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
                 "BME 129C Capstone  |  Progress Report",
                 font_size=16, color=GRAY)

    add_text_box(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.5),
                 "Sage Clokey   |   Spring 2026   |   UC Santa Cruz",
                 font_size=18, color=GREEN_LIGHT)

    # Add the three-layer diagram as a visual on the title slide
    if (FIGURES / "slide_three_layers.png").exists():
        slide.shapes.add_picture(
            str(FIGURES / "slide_three_layers.png"),
            Inches(7.5), Inches(4.8), Inches(5.5)
        )

    # =========================================================================
    # SLIDE 2: Goal & Approach (0.5 min)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
                 "Goal & Approach", font_size=32, color=GOLD, bold=True)

    bullets = [
        ("Are living systems literally decentralized economies?", [
            "Not metaphor -- measurable structural equivalence"
        ]),
        ("Thesis: biological networks exhibit the same properties Austrian economists attribute to free markets", [
            "Scale-free topology, distributed information flow, robustness to failure"
        ]),
        ("Three layers of convergent evidence from real data", [
            "Layer 1: Network topology (STRING DB, RegulonDB, KEGG)",
            "Layer 2: Economic simulation (pathways as agents vs. central allocator)",
            "Layer 3: Cross-species trade (comparative advantage via codon distance)",
        ]),
    ]
    add_bullet_slide_content(slide, bullets,
                              Inches(0.8), Inches(1.2), Inches(5.5), Inches(5.5),
                              font_size=16)

    # Add three-layer diagram on the right
    if (FIGURES / "slide_three_layers.png").exists():
        slide.shapes.add_picture(
            str(FIGURES / "slide_three_layers.png"),
            Inches(6.8), Inches(1.5), Inches(6.2)
        )

    # =========================================================================
    # SLIDE 3: What Centralized vs Decentralized Looks Like (visual)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "What Does Network Architecture Look Like?",
                 font_size=32, color=GOLD, bold=True)

    # Full-width network comparison image
    if (FIGURES / "slide_network_comparison.png").exists():
        slide.shapes.add_picture(
            str(FIGURES / "slide_network_comparison.png"),
            Inches(0.5), Inches(1.3), Inches(12.3)
        )

    # Labels below
    add_text_box(slide, Inches(0.5), Inches(5.8), Inches(3.8), Inches(1.2),
                 "One hub controls everything.\nRemove it = total collapse.",
                 font_size=14, color=RED)

    add_text_box(slide, Inches(4.5), Inches(5.8), Inches(4.5), Inches(1.2),
                 "Multiple hubs, high clustering.\nRobust to failure -- like a real economy.",
                 font_size=14, color=GREEN_MID)

    add_text_box(slide, Inches(9.2), Inches(5.8), Inches(3.8), Inches(1.2),
                 "No structure, no hubs.\nResilient but inefficient -- no specialization.",
                 font_size=14, color=GRAY)

    # =========================================================================
    # SLIDE 4: What's Built (1 min)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "Accomplished So Far", font_size=32, color=GOLD, bold=True)

    # Left column: Data
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
                 "Data Pipeline (Complete)", font_size=20, color=GREEN_MID, bold=True)

    data_bullets = [
        ("STRING PPI: E. coli (529 proteins, 6,951 interactions)", []),
        ("STRING PPI: Yeast (573 proteins, 6,342 interactions)", []),
        ("RegulonDB GRN: 282 nodes, 308 regulatory edges", []),
        ("KEGG metabolic networks: E. coli + yeast", []),
        ("Single-cell RNA-seq: PBMC 3k (scanpy)", []),
        ("All data cached locally for reproducibility", []),
    ]
    add_bullet_slide_content(slide, data_bullets,
                              Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5),
                              font_size=14, color=OFF_WHITE)

    # Right column: Analysis
    add_text_box(slide, Inches(7), Inches(1.2), Inches(5.5), Inches(0.5),
                 "Analysis Framework (Complete)", font_size=20, color=GREEN_MID, bold=True)

    analysis_bullets = [
        ("Topology: degree distribution, power-law fitting, Gini coefficient", []),
        ("Robustness curves: random failure vs. targeted attack", []),
        ("5 centralized reference networks for comparison", []),
        ("Economic simulation: distributed vs. centralized allocation", []),
        ("Cross-species trade cost matrix (codon distance)", []),
        ("Automated figure generation pipeline", []),
    ]
    add_bullet_slide_content(slide, analysis_bullets,
                              Inches(7), Inches(1.8), Inches(5.8), Inches(4.5),
                              font_size=14, color=OFF_WHITE)

    # =========================================================================
    # SLIDE 5: Figure 2 — The Data (1.5 min)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 "Figure 2: Biological Networks Are Decentralized",
                 font_size=28, color=GOLD, bold=True)

    # Main figure
    if (FIGURES / "figure2_network_topology.png").exists():
        slide.shapes.add_picture(
            str(FIGURES / "figure2_network_topology.png"),
            Inches(0.3), Inches(1.0), Inches(12.7), Inches(4.2)
        )

    # Caption
    add_text_box(slide, Inches(0.5), Inches(5.4), Inches(12), Inches(1.8),
                 "Figure 2. Topological analysis of E. coli and yeast PPI networks (STRING DB) "
                 "and E. coli GRN (RegulonDB) compared to centralized references. "
                 "(A) Power-law degree distributions (alpha = 2.05-2.44). "
                 "(B) Betweenness Gini shows hubs exist but are distributed across multiple nodes. "
                 "(C) PPI networks survive removing ~37% of top hubs; star/hub-spoke collapse at ~2%.",
                 font_size=12, color=GRAY)

    # =========================================================================
    # SLIDE 6: Robustness — The Key Insight
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "The Key Insight: Hubs Exist, But the System Doesn't Depend on Them",
                 font_size=26, color=GOLD, bold=True)

    # Robustness visual
    if (FIGURES / "slide_robustness_detail.png").exists():
        slide.shapes.add_picture(
            str(FIGURES / "slide_robustness_detail.png"),
            Inches(0.5), Inches(1.2), Inches(7.5), Inches(5.5)
        )

    # Key numbers on the right
    add_text_box(slide, Inches(8.5), Inches(1.5), Inches(4.5), Inches(0.5),
                 "Robustness Comparison", font_size=20, color=GREEN_MID, bold=True)

    # Data table as text
    table_text = (
        "Nodes removed before collapse:\n\n"
        "E. coli PPI (targeted):     37%\n"
        "Yeast PPI (targeted):       37%\n"
        "Star graph (targeted):        2%\n"
        "Hub-and-spoke (targeted):  2%\n\n"
        "Biological networks have\n"
        "multiple redundant paths.\n\n"
        "Centralized architectures\n"
        "have a single point of failure."
    )
    add_text_box(slide, Inches(8.5), Inches(2.2), Inches(4.3), Inches(4.5),
                 table_text, font_size=15, color=OFF_WHITE)

    # Highlight box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.5), Inches(5.8), Inches(4.3), Inches(1.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1b, 0x43, 0x32)
    shape.line.color.rgb = GREEN_MID
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Life is not a machine with a master switch.\nIt's an economy of cooperating agents."
    p.font.size = Pt(13)
    p.font.color.rgb = GREEN_LIGHT
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 7: Key Numbers Table
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "Summary of Results", font_size=32, color=GOLD, bold=True)

    # Build a proper table
    rows, cols = 6, 5
    tbl_shape = slide.shapes.add_table(rows, cols,
                                        Inches(0.8), Inches(1.3),
                                        Inches(11.5), Inches(4.5))
    tbl = tbl_shape.table

    headers = ["Network", "Nodes / Edges", "Power-law alpha", "Clustering", "Robustness (targeted)"]
    data_rows = [
        ["E. coli PPI", "529 / 6,951", "2.05", "0.645", "36.8%"],
        ["Yeast PPI", "573 / 6,342", "2.44", "0.651", "36.8%"],
        ["E. coli GRN", "282 / 308", "2.41", "0.000", "1.9%"],
        ["Star (centralized)", "529 / 528", "0.00", "0.000", "1.9%"],
        ["Hub-and-spoke", "529 / 564", "1.25", "0.000", "1.9%"],
    ]

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(14)
            para.font.bold = True
            para.font.color.rgb = GOLD
            para.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1b, 0x43, 0x32)

    for i, row_data in enumerate(data_rows):
        is_bio = i < 3
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(13)
                para.font.color.rgb = GREEN_LIGHT if is_bio else RED
                para.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL_BG

    # Interpretation below table
    add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(1),
                 "Biological networks: scale-free degree distributions (alpha 2-3), "
                 "high clustering (local cooperation), and 18x more robust than centralized architectures under targeted attack.",
                 font_size=15, color=GREEN_MID)

    # =========================================================================
    # SLIDE 8: Next Steps
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "Next Steps", font_size=32, color=GOLD, bold=True)

    next_bullets = [
        ("Layer 2: Distributed vs. centralized metabolic simulation", [
            "Does the 'market' outperform the 'planner' under perturbation?"
        ]),
        ("Motif analysis: classify network motifs", [
            "Entrepreneurial = local feedback loops, adaptive",
            "Central planning = hub-dependent cascades, brittle"
        ]),
        ("Layer 3: Cross-species trade network", [
            "Does codon distance predict gene transfer cost?",
            "Like the gravity model predicts international trade?"
        ]),
        ("Paper: Results and discussion for all three layers", []),
    ]
    add_bullet_slide_content(slide, next_bullets,
                              Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5),
                              font_size=17)

    # Add layer 2 economy figure preview on right
    if (FIGURES / "layer2_economy.png").exists():
        slide.shapes.add_picture(
            str(FIGURES / "layer2_economy.png"),
            Inches(7), Inches(1.5), Inches(5.8), Inches(4.0)
        )
        add_text_box(slide, Inches(7), Inches(5.6), Inches(5.8), Inches(0.5),
                     "Preview: Layer 2 economic simulation results",
                     font_size=12, color=GRAY)

    # Save
    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
