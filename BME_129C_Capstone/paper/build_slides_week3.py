"""
Build PowerPoint progress report slides for BME 129C Capstone — Week 3.
3 minutes: 0.5 min goal/approach, 1 min accomplishments, 1.5 min recent progress.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

FIGURES = Path(__file__).resolve().parent / "figures"
OUT = Path(__file__).resolve().parent / "progress_report_slides_week3.pptx"

# Colors
GREEN_DARK = RGBColor(0x2d, 0x6a, 0x4f)
GREEN_MID = RGBColor(0x52, 0xb7, 0x88)
GREEN_LIGHT = RGBColor(0x95, 0xd5, 0xb2)
GOLD = RGBColor(0xe9, 0xc4, 0x6a)
RED = RGBColor(0xe6, 0x39, 0x46)
WHITE = RGBColor(0xff, 0xff, 0xff)
OFF_WHITE = RGBColor(0xf5, 0xf5, 0xf5)
DARK_BG = RGBColor(0x0d, 0x11, 0x17)
PANEL_BG = RGBColor(0x16, 0x1b, 0x22)
GRAY = RGBColor(0x8b, 0x94, 0x9e)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    return txBox


def add_bullets(slide, bullets, left, top, width, height,
                font_size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, subs) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
        for sub in subs:
            p2 = tf.add_paragraph()
            p2.text = sub
            p2.font.size = Pt(font_size - 2)
            p2.font.color.rgb = GRAY
            p2.font.name = "Calibri"
            p2.space_after = Pt(4)
            p2.level = 1
    return txBox


def add_green_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN_MID
    shape.line.fill.background()


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # =========================================================================
    # SLIDE 1: Title (part of 0.5 min)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                 "Living Systems as Decentralized Economies",
                 font_size=40, color=GREEN_MID, bold=True)
    add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
                 "Why the Molecular Biologist Should Be an Entrepreneur, Not a Central Planner",
                 font_size=22, color=GOLD)
    add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
                 "BME 129C Capstone  |  Week 3 Progress Report",
                 font_size=16, color=GRAY)
    add_text_box(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.5),
                 "Sage Clokey   |   Spring 2026   |   UC Santa Cruz",
                 font_size=18, color=GREEN_LIGHT)

    if (FIGURES / "slide_three_layers.png").exists():
        slide.shapes.add_picture(
            str(FIGURES / "slide_three_layers.png"),
            Inches(7.5), Inches(4.8), Inches(5.5))

    # =========================================================================
    # SLIDE 2: Goal & Approach (0.5 min)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
                 "Goal & Approach", font_size=32, color=GOLD, bold=True)

    bullets = [
        ("Question: Are living systems literally decentralized economies?", [
            "Not metaphor — measurable structural equivalence"
        ]),
        ("Thesis: biological networks exhibit the same properties Austrian economists attribute to free markets", [
            "Scale-free topology, distributed information flow, robustness to failure"
        ]),
        ("Three layers of convergent evidence from real data:", [
            "Layer 1: Network topology — no master node (RegulonDB, STRING, KEGG)",
            "Layer 1b: Single-cell economy — cells as specialized agents (10x Genomics pbmc3k)",
            "Layer 2: Economic simulation — distributed vs. centralized allocation",
            "Layer 3: Cross-species trade — comparative advantage via codon distance",
        ]),
    ]
    add_bullets(slide, bullets, Inches(0.8), Inches(1.2), Inches(5.5), Inches(5.5), font_size=16)

    if (FIGURES / "slide_three_layers.png").exists():
        slide.shapes.add_picture(
            str(FIGURES / "slide_three_layers.png"),
            Inches(6.8), Inches(1.5), Inches(6.2))

    # =========================================================================
    # SLIDE 3: Accomplishments — All Layers Complete (1 min)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "Accomplished So Far — All Layers Complete",
                 font_size=32, color=GOLD, bold=True)

    # Left column
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
                 "Data & Analysis Pipeline", font_size=20, color=GREEN_MID, bold=True)

    left = [
        ("Layer 1: Gene regulatory, metabolic, and PPI networks", [
            "E. coli PPI: 529 nodes, 6,951 edges",
            "Yeast PPI: 573 nodes, 6,342 edges",
            "RegulonDB GRN: 282 nodes, 308 edges",
            "Power-law exponents: α = 2.05–2.44"
        ]),
        ("Layer 1b: Single-cell economy (PBMC 3k)", [
            "8 cell types, no master cell",
            "Communication Gini = 0.0 (fully distributed)",
            "75% communication survival after any cell type removal"
        ]),
    ]
    add_bullets(slide, left, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), font_size=13, color=OFF_WHITE)

    # Right column
    add_text_box(slide, Inches(7), Inches(1.2), Inches(5.5), Inches(0.5),
                 "Simulation & Trade Analysis", font_size=20, color=GREEN_MID, bold=True)

    right = [
        ("Layer 2: Economic simulation (13 pathway agents)", [
            "Distributed: 2.9× lower variance than centralized",
            "71.1% GDP retained under perturbation (vs 53.0%)",
            "Local feedback alone reaches stable equilibrium"
        ]),
        ("Layer 3: Cross-species trade (8 organisms, 4 kingdoms)", [
            "Within-kingdom cost: 0.17–0.38",
            "Cross-kingdom cost: 0.65–0.83",
            "Each organism contributes unique comparative advantage"
        ]),
    ]
    add_bullets(slide, right, Inches(7), Inches(1.8), Inches(5.8), Inches(5.0), font_size=13, color=OFF_WHITE)

    # Bottom bar
    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
                 "Paper draft complete (Abstract → Results)  |  Figure guide with falsifiable predictions  |  Reproducible pipeline (run_all.py)",
                 font_size=13, color=GREEN_MID)

    # =========================================================================
    # SLIDE 4: Key Figure — Network Topology (1.5 min — part 1)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 "Layer 1: Biological Networks Are Decentralized",
                 font_size=28, color=GOLD, bold=True)

    if (FIGURES / "figure2_network_topology.png").exists():
        slide.shapes.add_picture(
            str(FIGURES / "figure2_network_topology.png"),
            Inches(0.3), Inches(1.0), Inches(12.7), Inches(4.2))

    add_text_box(slide, Inches(0.5), Inches(5.4), Inches(12), Inches(1.8),
                 "(A) Power-law degree distributions (α = 2.05–2.44) — many nodes with few connections, "
                 "few hubs, no single controller. "
                 "(B) Betweenness Gini: hubs exist but are distributed — no single bottleneck. "
                 "(C) Robustness: PPI networks survive ~37% targeted hub removal; "
                 "star/hub-spoke collapse at ~2%. A 19:1 robustness ratio.",
                 font_size=12, color=GRAY)

    # =========================================================================
    # SLIDE 5: Key Figure — Single-Cell + Layer 2 (1.5 min — part 2)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 "Layer 1b + Layer 2: Cells as Agents, Pathways as Markets",
                 font_size=28, color=GOLD, bold=True)

    # Single-cell figure on left
    if (FIGURES / "layer1b_single_cell_economy.png").exists():
        slide.shapes.add_picture(
            str(FIGURES / "layer1b_single_cell_economy.png"),
            Inches(0.3), Inches(1.0), Inches(6.2), Inches(4.5))

    # Layer 2 figure on right
    if (FIGURES / "layer2_economy.png").exists():
        slide.shapes.add_picture(
            str(FIGURES / "layer2_economy.png"),
            Inches(6.8), Inches(1.0), Inches(6.2), Inches(4.5))

    add_text_box(slide, Inches(0.3), Inches(5.7), Inches(6.2), Inches(1.5),
                 "Left: 2,638 cells, same genome, 8 specialized types.\n"
                 "No master cell. Communication Gini = 0.0.\n"
                 "The immune system is a decentralized economy.",
                 font_size=12, color=GRAY)

    add_text_box(slide, Inches(6.8), Inches(5.7), Inches(6.2), Inches(1.5),
                 "Right: Distributed allocation (green) vs. centralized (red).\n"
                 "Distributed reaches equilibrium through local feedback alone.\n"
                 "Retains 71% GDP under perturbation vs. 53% for the planner.",
                 font_size=12, color=GRAY)

    # =========================================================================
    # SLIDE 6: Key Figure — Trade Network (1.5 min — part 3)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
                 "Layer 3: The Tree of Life Is a Trade Network",
                 font_size=28, color=GOLD, bold=True)

    if (FIGURES / "layer3_trade_network.png").exists():
        slide.shapes.add_picture(
            str(FIGURES / "layer3_trade_network.png"),
            Inches(0.5), Inches(1.0), Inches(6.0), Inches(5.5))

    # Key points on right
    trade_bullets = [
        ("8 organisms across 4 kingdoms", []),
        ("Codon usage distance = trade friction", [
            "Within-kingdom: 0.17–0.38 (free trade zone)",
            "Cross-kingdom: 0.65–0.83 (trade embargo)"
        ]),
        ("Each organism has unique comparative advantage", [
            "Coral → biomineralization",
            "Spider → silk proteins",
            "Axolotl → regeneration",
            "Bacteria → cellulose synthesis"
        ]),
        ("The tree of life forms natural trade blocs", [
            "Like economic free trade zones structured by evolutionary distance"
        ]),
    ]
    add_bullets(slide, trade_bullets,
                Inches(7), Inches(1.2), Inches(5.8), Inches(5.5),
                font_size=15, color=OFF_WHITE)

    # =========================================================================
    # SLIDE 7: This Week's Progress + Advisors
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "This Week's Progress", font_size=32, color=GOLD, bold=True)

    progress = [
        ("Faculty advisors secured", [
            "Todd Lowe (RNA biology, genomics) — co-advisor, meeting Wed 4:30pm",
            "Josh Stuart (computational genomics) — connected through Dubois, meeting this week"
        ]),
        ("Capstone figures presentation updated (v2)", [
            "Refined visualizations for advisor meetings"
        ]),
        ("Methods section drafted", [
            "Arial 11, single-spaced, database table with accession numbers",
            "Data Availability section: GitHub + future Dryad deposit"
        ]),
        ("Paper argument sharpened through faculty exchanges", [
            "Bernick: 'programming is Plato's shadow puppets' — connects to machine metaphor critique",
            "Etymology of 'program' (pre-written) as evidence the field's language misrepresents living systems"
        ]),
    ]
    add_bullets(slide, progress, Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5), font_size=15)

    # =========================================================================
    # SLIDE 8: Next Steps
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_green_bar(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "Next Steps", font_size=32, color=GOLD, bold=True)

    nexts = [
        ("Validate Layer 2 with real flux data", [
            "Integrate COBRApy flux balance analysis from BiGG E. coli model",
            "Replace simplified parameterization with genome-scale metabolic data"
        ]),
        ("Meet with advisors Lowe and Stuart", [
            "Present figures and paper draft",
            "Get feedback on methodology and framing"
        ]),
        ("Complete Discussion section of paper", [
            "Disease as broken decentralization (cancer, drug resistance)",
            "Practical implications: design distributed feedback, not command hierarchies"
        ]),
        ("Statistical validation", [
            "Kolmogorov-Smirnov test for power-law fit",
            "Spearman correlation: codon distance vs. trade cost (gravity model)"
        ]),
    ]
    add_bullets(slide, nexts, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5), font_size=16)

    # Highlight box on right
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.5), Inches(1.5), Inches(5), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1b, 0x43, 0x32)
    shape.line.color.rgb = GREEN_MID
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The Thesis in One Sentence"
    p.font.size = Pt(18)
    p.font.color.rgb = GOLD
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = (
        "\nLife organizes as a decentralized economy — "
        "not a machine with a master switch, "
        "but a network of cooperating agents "
        "coordinating through local signals.\n\n"
        "The data shows it at every scale:\n"
        "genes, cells, pathways, species.\n\n"
        "The molecular biologist should be\n"
        "an entrepreneur, not a central planner."
    )
    p2.font.size = Pt(14)
    p2.font.color.rgb = GREEN_LIGHT
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
