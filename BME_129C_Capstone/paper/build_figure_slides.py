"""
Build PowerPoint: Figure-by-figure walkthrough of capstone results.
Each figure gets its own slide with interpretation in the Sagent Creed voice.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FIGURES = Path(__file__).resolve().parent / "figures"
OUT = Path(__file__).resolve().parent / "capstone_figures_v2.pptx"

# Colors
GREEN_DARK  = RGBColor(0x1b, 0x43, 0x32)
GREEN_MID   = RGBColor(0x2d, 0x6a, 0x4f)
GREEN_HI    = RGBColor(0x52, 0xb7, 0x88)
GREEN_LIGHT = RGBColor(0x95, 0xd5, 0xb2)
GOLD        = RGBColor(0xe9, 0xc4, 0x6a)
RED         = RGBColor(0xe6, 0x39, 0x46)
WHITE       = RGBColor(0xff, 0xff, 0xff)
OFF_WHITE   = RGBColor(0xf0, 0xeb, 0xe0)
DARK_BG     = RGBColor(0x06, 0x0b, 0x06)
PANEL_BG    = RGBColor(0x0f, 0x1a, 0x0f)
GRAY        = RGBColor(0x8b, 0x94, 0x9e)
DIM         = RGBColor(0x5a, 0x6a, 0x56)

SW = Inches(13.33)
SH = Inches(7.5)


def set_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def green_bar(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.05))
    s.fill.solid()
    s.fill.fore_color.rgb = GREEN_HI
    s.line.fill.background()


def gold_bar(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.42), SW, Inches(0.08))
    s.fill.solid()
    s.fill.fore_color.rgb = GOLD
    s.line.fill.background()


def add_text(slide, left, top, w, h, text, size=18, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, font="Georgia", spacing=None):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    if spacing:
        p.space_after = spacing
    return box


def add_multi(slide, left, top, w, h, lines, size=16, color=WHITE, font="Georgia"):
    """Add text box with multiple paragraphs from a list of (text, color, bold, size_override) tuples."""
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, c, b, s = item, color, False, size
        else:
            txt = item[0]
            c = item[1] if len(item) > 1 else color
            b = item[2] if len(item) > 2 else False
            s = item[3] if len(item) > 3 else size
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(s)
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = font
        p.space_after = Pt(6)
    return box


def add_img(slide, path, left, top, w=None, h=None):
    if not path.exists():
        return
    if w and h:
        slide.shapes.add_picture(str(path), left, top, w, h)
    elif w:
        slide.shapes.add_picture(str(path), left, top, width=w)
    elif h:
        slide.shapes.add_picture(str(path), left, top, height=h)
    else:
        slide.shapes.add_picture(str(path), left, top)


def quote_box(slide, left, top, w, h, text, cite=""):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = PANEL_BG
    s.line.color.rgb = GREEN_MID
    s.line.width = Pt(1)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = GOLD
    p.font.name = "Georgia"
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER
    if cite:
        p2 = tf.add_paragraph()
        p2.text = cite
        p2.font.size = Pt(10)
        p2.font.color.rgb = DIM
        p2.font.name = "Georgia"
        p2.alignment = PP_ALIGN.CENTER


def build():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    # =====================================================================
    # SLIDE 1: Title
    # =====================================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    green_bar(sl)
    gold_bar(sl)

    add_text(sl, Inches(1.5), Inches(1.2), Inches(10), Inches(1.5),
             "Living Systems as\nDecentralized Economies",
             size=42, color=GREEN_HI, bold=True)

    add_text(sl, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
             "Why the Molecular Biologist Should Be a Sagent, Not a Central Planner",
             size=20, color=GOLD)

    quote_box(sl, Inches(1.5), Inches(4.2), Inches(6), Inches(1.2),
              "The first job was gardener, not king.\nThe mandate was to tend — not to replace.",
              "— Genesis 2:15")

    add_text(sl, Inches(1.5), Inches(6.0), Inches(10), Inches(0.5),
             "Sage Clokey  |  BME 129C  |  Spring 2026  |  UC Santa Cruz",
             size=14, color=DIM)

    add_img(sl, FIGURES / "slide_three_layers.png",
            Inches(8), Inches(3.8), w=Inches(4.8))

    # =====================================================================
    # SLIDE 2: The Question
    # =====================================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    green_bar(sl)

    add_text(sl, Inches(1), Inches(0.4), Inches(11), Inches(0.5),
             "The Question", size=30, color=GOLD, bold=True)

    add_multi(sl, Inches(1), Inches(1.2), Inches(5.5), Inches(5.5), [
        ("Where life grows in spirals,\nwe imposed straight lines.", GOLD, True, 22),
        ("", WHITE, False, 8),
        ("Are living systems literally decentralized economies?", WHITE, False, 18),
        ("Not metaphor — measurable structural equivalence.", GREEN_LIGHT, False, 16),
        ("", WHITE, False, 8),
        ("The block: rigid, centralized, imposed from outside.", RED, False, 16),
        ("The spiral: adaptive, distributed, grown from within.", GREEN_HI, False, 16),
        ("", WHITE, False, 8),
        ("Three layers of evidence from real biological data:", WHITE, False, 16),
        ("Layer 1: Network topology (STRING DB, RegulonDB)", DIM, False, 14),
        ("Layer 2: Economic simulation (pathways as agents)", DIM, False, 14),
        ("Layer 3: Cross-species trade (codon distance)", DIM, False, 14),
    ])

    add_img(sl, FIGURES / "slide_network_comparison.png",
            Inches(6.8), Inches(1.5), w=Inches(6))

    add_multi(sl, Inches(7), Inches(5.3), Inches(5.5), Inches(1.5), [
        ("Centralized (Star)", RED, False, 11),
        ("Biological (Scale-Free)", GREEN_HI, False, 11),
        ("Random (No Structure)", GRAY, False, 11),
    ])

    # =====================================================================
    # SLIDE 2b: The Austrian Economists
    # =====================================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    green_bar(sl)

    add_text(sl, Inches(0.6), Inches(0.2), Inches(12), Inches(0.5),
             "The Austrian Economists Who Predicted the Data",
             size=28, color=GOLD, bold=True)

    add_text(sl, Inches(0.6), Inches(0.7), Inches(12), Inches(0.3),
             "Five predictions from economics — five confirmations from molecular biology.",
             size=14, color=DIM, font="Georgia")

    # Economist portraits and descriptions — 5 across
    econ_data = [
        ("economist_menger.png", "Carl Menger\n(1840–1921)",
         "Spontaneous order\narises from individual\naction — no designer\nneeded at the top."),
        ("economist_mises.png", "Ludwig von Mises\n(1881–1973)",
         "The calculation problem:\ncentral planners cannot\nreplicate the information\nthat prices discover."),
        ("economist_hayek.png", "Friedrich Hayek\n(1899–1992)",
         "Knowledge is dispersed.\nNo single node can hold\nwhat the whole network\nknows collectively."),
        ("economist_rothbard.png", "Murray Rothbard\n(1926–1995)",
         "Centralization creates\na single point of failure.\nConcentrate authority and\nyou guarantee collapse."),
        ("economist_kirzner.png", "Israel Kirzner\n(b. 1930)",
         "Entrepreneurial discovery:\nagents find opportunities\nthrough local feedback,\nnot central assignment."),
    ]

    for i, (img, name, desc) in enumerate(econ_data):
        x = Inches(0.4 + i * 2.6)
        add_img(sl, FIGURES / img, x, Inches(1.2), w=Inches(2.0), h=Inches(2.0))
        add_text(sl, x, Inches(3.3), Inches(2.2), Inches(0.6),
                 name, size=12, color=GREEN_HI, bold=True, align=PP_ALIGN.CENTER)
        add_text(sl, x, Inches(4.0), Inches(2.3), Inches(1.8),
                 desc, size=10, color=OFF_WHITE, align=PP_ALIGN.CENTER)

    quote_box(sl, Inches(1), Inches(5.8), Inches(11), Inches(1.2),
              "These men described how free economies work.\n"
              "Molecular biology proves their predictions — not in markets, but in cells.",
              "")

    # =====================================================================
    # SLIDE 3: Figure 1 — Network Topology (full)
    # =====================================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    green_bar(sl)

    add_text(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5),
             "Figure 1: No Master Node — Hayek's Knowledge Problem in the Cell",
             size=24, color=GOLD, bold=True)

    add_img(sl, FIGURES / "layer1_topology.png",
            Inches(0.3), Inches(0.9), w=Inches(12.7), h=Inches(4.2))

    # Hayek portrait
    add_img(sl, FIGURES / "economist_hayek.png",
            Inches(11.3), Inches(5.3), w=Inches(1.6), h=Inches(1.6))
    add_text(sl, Inches(11.0), Inches(6.9), Inches(2.2), Inches(0.4),
             "F. A. Hayek", size=9, color=DIM, align=PP_ALIGN.CENTER)

    add_multi(sl, Inches(0.5), Inches(5.3), Inches(10.5), Inches(2), [
        ("Hayek (1945): \"Knowledge is dispersed among many minds. No central authority can aggregate it.\"", GOLD, True, 13),
        ("Biological networks follow power-law degree distributions (α = 2.05–2.44) — many nodes with few connections, a few hubs, but no single node that dominates.", GREEN_LIGHT, False, 12),
        ("Betweenness Gini — hubs exist in biology (0.85–0.94) but the star graph (1.00) routes ALL paths through one node.", WHITE, False, 12),
        ("PPI networks survive removing ~37–49% of nodes. Star graph collapses at 2%. Hayek was right: the network knows more than any node.", WHITE, True, 13),
    ])

    # =====================================================================
    # SLIDE 4: Figure 2 — PPI Focused Comparison
    # =====================================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    green_bar(sl)

    add_text(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5),
             "Figure 2: Rothbard's Warning — The Block Measured Against the Spiral",
             size=24, color=GOLD, bold=True)

    add_img(sl, FIGURES / "figure2_network_topology.png",
            Inches(0.3), Inches(0.9), w=Inches(12.7), h=Inches(4.0))

    add_multi(sl, Inches(0.5), Inches(5.1), Inches(6), Inches(2.2), [
        ("The key comparison:", GOLD, True, 16),
        ("", WHITE, False, 4),
        ("E. coli PPI: 529 proteins, 6,951 interactions", GREEN_LIGHT, False, 14),
        ("  Power-law α = 2.05  |  Clustering = 0.645", GREEN_LIGHT, False, 13),
        ("  Survives removing 37% of top hubs", GREEN_LIGHT, False, 13),
        ("", WHITE, False, 4),
        ("Star graph: 529 nodes, 528 edges", RED, False, 14),
        ("  One hub controls everything", RED, False, 13),
        ("  Collapses at 2% removal", RED, False, 13),
    ])

    # Rothbard portrait
    add_img(sl, FIGURES / "economist_rothbard.png",
            Inches(7.2), Inches(5.2), w=Inches(1.4), h=Inches(1.4))

    # Interpretation box on right
    quote_box(sl, Inches(8.8), Inches(5.1), Inches(4.2), Inches(2.0),
              "19:1 robustness ratio.\n\n"
              "Rothbard: \"Concentrate authority\n"
              "in one node and you guarantee\n"
              "catastrophic failure.\"\n\n"
              "The spiral was designed to endure.",
              "— M. Rothbard")

    # =====================================================================
    # SLIDE 5: Robustness Visual
    # =====================================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    green_bar(sl)

    add_text(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5),
             "What Happens When You Attack the Hubs?",
             size=24, color=GOLD, bold=True)

    add_img(sl, FIGURES / "slide_robustness_detail.png",
            Inches(0.4), Inches(1.0), w=Inches(8), h=Inches(5.8))

    add_multi(sl, Inches(8.8), Inches(1.2), Inches(4), Inches(5.5), [
        ("Hubs exist in biology.", WHITE, True, 18),
        ("But the system doesn't\ndepend on them.", GREEN_HI, True, 18),
        ("", WHITE, False, 10),
        ("Remove the top 5 hubs from a\nbiological network:", WHITE, False, 15),
        ("Still 82% connected.", GREEN_HI, True, 20),
        ("", WHITE, False, 10),
        ("Remove the single hub from\na star graph:", WHITE, False, 15),
        ("0% connected.", RED, True, 20),
        ("Total collapse.", RED, False, 15),
        ("", WHITE, False, 10),
        ("This is Rothbard's warning:\nconcentrate authority in one node,\nand you create a single point\nof catastrophic failure.", DIM, False, 13),
    ])

    # Rothbard portrait bottom-right
    add_img(sl, FIGURES / "economist_rothbard.png",
            Inches(11.0), Inches(5.5), w=Inches(1.5), h=Inches(1.5))
    add_text(sl, Inches(10.7), Inches(7.0), Inches(2.2), Inches(0.4),
             "M. Rothbard", size=9, color=DIM, align=PP_ALIGN.CENTER)

    # =====================================================================
    # SLIDE 6: Layer 1b — Single-Cell Economy
    # =====================================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    green_bar(sl)

    add_text(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5),
             "Figure 3: Cells as Sagents — Same Genome, Different Callings",
             size=24, color=GOLD, bold=True)

    add_img(sl, FIGURES / "layer1b_single_cell_economy.png",
            Inches(0.3), Inches(0.9), w=Inches(12.7), h=Inches(4.0))

    add_multi(sl, Inches(0.5), Inches(5.1), Inches(6), Inches(2.2), [
        ("37 trillion cells. One genome. No master cell.", GOLD, True, 16),
        ("", WHITE, False, 4),
        ("Left: UMAP shows 8 distinct cell types — each a 'sector'\nof the cellular economy, producing different goods\nfrom the same raw material (genome).", WHITE, False, 13),
        ("Center: Division of labor — CD4 T cells are specialists,\nmegakaryocytes are generalists. Same DNA, different output.", WHITE, False, 13),
        ("Right: Communication network — Betweenness Gini = 0.000.\nNo gatekeeper. Every cell type talks directly to every other.", GREEN_LIGHT, False, 13),
    ])

    # Menger portrait
    add_img(sl, FIGURES / "economist_menger.png",
            Inches(7.0), Inches(5.2), w=Inches(1.4), h=Inches(1.4))

    quote_box(sl, Inches(8.6), Inches(5.1), Inches(4.2), Inches(2.0),
              "Menger's spontaneous order\nmade visible: complex differentiation\narising from individual cells\nresponding to local conditions.\n\n"
              "No cell was told what to become.",
              "— C. Menger, Principles (1871)")

    # =====================================================================
    # SLIDE 7: Layer 2 — Economic Modeling
    # =====================================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    green_bar(sl)

    add_text(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5),
             "Figure 4: Mises' Calculation Problem — The Planner's Plan Breaks",
             size=24, color=GOLD, bold=True)

    add_text(sl, Inches(0.6), Inches(0.75), Inches(12), Inches(0.3),
             "The block looks efficient on paper. It always does — until conditions change.",
             size=14, color=DIM, font="Georgia")

    add_img(sl, FIGURES / "layer2_economy.png",
            Inches(0.3), Inches(1.1), w=Inches(12.7), h=Inches(3.8))

    # Mises portrait — left panel
    add_img(sl, FIGURES / "economist_mises.png",
            Inches(0.5), Inches(5.1), w=Inches(1.1), h=Inches(1.1))

    add_multi(sl, Inches(1.7), Inches(5.1), Inches(2.8), Inches(2.2), [
        ("GDP Over Time", GOLD, True, 14),
        ("Centralized (red) reaches higher output —\nbut 2.9x higher variance.", WHITE, False, 12),
        ("Mises (1920): central planners cannot\ncalculate without real prices.", GOLD, False, 11),
        ("Distributed (green) finds equilibrium\nthrough local feedback alone.", GREEN_LIGHT, False, 12),
    ])

    add_multi(sl, Inches(4.7), Inches(5.1), Inches(4), Inches(2.2), [
        ("Robustness: Mises Proven Right", GOLD, True, 14),
        ("Remove HIF1α pathway:", WHITE, False, 12),
        ("  Distributed retains 71.1% GDP", GREEN_HI, True, 14),
        ("  Centralized retains 53.0% GDP", RED, True, 14),
        ("18.1 percentage point gap.", WHITE, False, 12),
        ("The calculation problem, measured.", GOLD, False, 12),
    ])

    # Kirzner portrait — right panel
    add_img(sl, FIGURES / "economist_kirzner.png",
            Inches(11.8), Inches(5.1), w=Inches(1.1), h=Inches(1.1))

    add_multi(sl, Inches(9), Inches(5.1), Inches(2.7), Inches(2.2), [
        ("Kirzner's Discovery", GOLD, True, 14),
        ("Agents start at uniform rates.\nThrough feedback, each discovers\nits optimal production level.", WHITE, False, 12),
        ("The oscillations are not a bug —\nthey ARE the entrepreneurial\ndiscovery process.", GREEN_LIGHT, False, 12),
    ])

    add_text(sl, Inches(1.7), Inches(6.3), Inches(2), Inches(0.3),
             "L. von Mises", size=9, color=DIM, align=PP_ALIGN.CENTER)
    add_text(sl, Inches(11.5), Inches(6.3), Inches(2), Inches(0.3),
             "I. Kirzner", size=9, color=DIM, align=PP_ALIGN.CENTER)

    # =====================================================================
    # SLIDE 8: Layer 3 — Trade Network
    # =====================================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    green_bar(sl)

    add_text(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5),
             "Figure 5: Menger's Spontaneous Order Across Kingdoms",
             size=24, color=GOLD, bold=True)

    add_text(sl, Inches(0.6), Inches(0.75), Inches(12), Inches(0.3),
             "No organism does everything. Each was given something the others lack.",
             size=14, color=DIM, font="Georgia")

    # Menger portrait top-right
    add_img(sl, FIGURES / "economist_menger.png",
            Inches(11.5), Inches(0.15), w=Inches(1.2), h=Inches(1.2))

    add_img(sl, FIGURES / "layer3_trade_network.png",
            Inches(0.3), Inches(1.1), w=Inches(12.7), h=Inches(3.8))

    add_multi(sl, Inches(0.5), Inches(5.1), Inches(4), Inches(2.2), [
        ("Left: Trade Network", GOLD, True, 14),
        ("Thick edges = easy trade\n(same kingdom, compatible codons).", WHITE, False, 12),
        ("Thin edges = high friction\n(cross-kingdom, different regulatory\nmachinery).", WHITE, False, 12),
        ("No biological hegemon.\nA distributed network of specialists.", GREEN_LIGHT, False, 12),
    ])

    add_multi(sl, Inches(4.7), Inches(5.1), Inches(4), Inches(2.2), [
        ("Center: Trade Cost Heatmap", GOLD, True, 14),
        ("Within-kingdom: 0.17–0.38", GREEN_HI, False, 13),
        ("Cross-kingdom: 0.65–0.83", RED, False, 13),
        ("The kingdom boundary is the\nsharpest jump in trade cost —", WHITE, False, 12),
        ("like trading across incompatible\nlegal systems.", GOLD, False, 12),
    ])

    add_multi(sl, Inches(9), Inches(5.1), Inches(4), Inches(2.2), [
        ("Right: Comparative Advantage", GOLD, True, 14),
        ("Coral → biomineralization", GREEN_LIGHT, False, 12),
        ("Spider → silk", GREEN_LIGHT, False, 12),
        ("Bacteria → cellulose", GREEN_LIGHT, False, 12),
        ("Axolotl → regeneration", GREEN_LIGHT, False, 12),
        ("The tree of life is a network\nof specialists, not generalists.", WHITE, False, 12),
        ("Trade is beneficial because\norganisms are different.", GOLD, False, 12),
    ])

    # =====================================================================
    # SLIDE 9: Summary Table
    # =====================================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    green_bar(sl)

    add_text(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5),
             "The Spiral Defeats the Block — At Every Scale",
             size=28, color=GOLD, bold=True)

    # Build summary table
    rows, cols = 8, 3
    tbl = sl.shapes.add_table(rows, cols,
                               Inches(0.8), Inches(1.2),
                               Inches(11.5), Inches(4.8)).table

    headers = ["Austrian Prediction", "Biological Evidence", "The Number"]
    data = [
        ("Knowledge is dispersed (Hayek)",
         "No master node in GRN, PPI, or metabolic networks",
         "Power-law α = 2.05–2.44"),
        ("Centralization creates fragility (Rothbard)",
         "Star/hub-spoke collapse under targeted attack",
         "19:1 robustness ratio"),
        ("Spontaneous order arises from individual action (Menger)",
         "Cell type specialization from identical genomes",
         "Gini = 0.000 (no gatekeeper)"),
        ("Central planning fails under change (Mises)",
         "Centralized allocation fails after perturbation",
         "71.1% vs 53.0% GDP retained"),
        ("Competitive discovery is the mechanism (Kirzner)",
         "Production rates converge through local feedback",
         "2.9x lower variance"),
        ("Trade follows institutional compatibility",
         "Codon distance predicts gene transfer cost",
         "0.17 within vs 0.83 across kingdoms"),
        ("Voluntary cooperation, not coerced obedience",
         "Cell communication: no gatekeeper, 75% fault-tolerant",
         "8 cell types, all connected"),
    ]

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = GOLD
            p.font.name = "Georgia"
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN_DARK

    for i, (c1, c2, c3) in enumerate(data):
        for j, val in enumerate((c1, c2, c3)):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = [OFF_WHITE, GREEN_LIGHT, GOLD][j]
                p.font.name = "Georgia"
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL_BG

    quote_box(sl, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.9),
              "Empire requires blocks. Freedom requires spirals. The data shows which one creation was built on.",
              "")

    # =====================================================================
    # SLIDE 10: Conclusion
    # =====================================================================
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    green_bar(sl)
    gold_bar(sl)

    add_text(sl, Inches(1.5), Inches(1.0), Inches(10), Inches(1),
             "Life is not a block you\ncarve into shape.",
             size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(sl, Inches(1.5), Inches(2.5), Inches(10), Inches(1),
             "It is a spiral you cultivate.",
             size=36, color=GREEN_HI, bold=True, align=PP_ALIGN.CENTER)

    add_multi(sl, Inches(2), Inches(3.8), Inches(9), Inches(2.5), [
        ("The sagent outperforms the central planner.", GOLD, True, 20),
        ("Not sometimes. Not under special conditions.", WHITE, False, 16),
        ("Structurally. At every scale. In every kingdom.", WHITE, False, 16),
        ("Over four billion years.", WHITE, False, 16),
        ("", WHITE, False, 10),
        ("The first job was gardener.", GREEN_LIGHT, True, 22),
        ("It still is.", GREEN_LIGHT, True, 22),
    ], font="Georgia")

    # Row of economist portraits at bottom
    econ_imgs = ["economist_menger.png", "economist_mises.png",
                 "economist_hayek.png", "economist_rothbard.png",
                 "economist_kirzner.png"]
    for i, img in enumerate(econ_imgs):
        add_img(sl, FIGURES / img,
                Inches(2.5 + i * 1.8), Inches(5.8), w=Inches(0.9), h=Inches(0.9))

    add_text(sl, Inches(1.5), Inches(6.8), Inches(10), Inches(0.5),
             "Sage Clokey  |  BME 129C  |  The Living Age",
             size=13, color=DIM, align=PP_ALIGN.CENTER)

    # Save
    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
