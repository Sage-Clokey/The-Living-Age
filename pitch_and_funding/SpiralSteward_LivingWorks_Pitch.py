import os, io, urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = "/mnt/c/Users/SajcS/Desktop/SpiralSteward_LivingWorks_Pitch.pptx"

# Colors
BG   = RGBColor(0x0d,0x11,0x17)
CARD = RGBColor(0x16,0x1b,0x22)
ELEV = RGBColor(0x1c,0x23,0x33)
GOLD = RGBColor(0xc9,0xa8,0x4c)
TEAL = RGBColor(0x4e,0xcd,0xc4)
GRN  = RGBColor(0x2d,0x6a,0x4f)
WHT  = RGBColor(0xe6,0xed,0xf3)
MUT  = RGBColor(0x8b,0x94,0x9e)
DIM  = RGBColor(0x48,0x4f,0x58)

BASE = "https://raw.githubusercontent.com/Sage-Clokey/Living-works-by-the-word/main/"
URLS = {
    "spiral_house":    BASE + "media/graphics/architecture/01_fibonacci_spiral_house.png",
    "mycelium_plan":   BASE + "media/graphics/architecture/02_mycelium_network_plan.png",
    "living_facade":   BASE + "media/graphics/architecture/03_living_facade.png",
    "branching":       BASE + "media/graphics/architecture/04_branching_tree_structure.png",
    "sunflower":       BASE + "media/graphics/architecture/05_sunflower_room_spiral.png",
    "block_vs_spiral": BASE + "media/graphics/architecture/06_block_vs_spiral_manifesto.png",
    "mycelium_dome":   BASE + "media/graphics/architecture/07_mycelium_dome_3d.png",
    "cg1": BASE + "media/graphics/architecture/ChatGPT%20Image%20Feb%2028%2C%202026%2C%2002_33_33%20PM.png",
    "cg2": BASE + "media/graphics/architecture/ChatGPT%20Image%20Feb%2028%2C%202026%2C%2002_35_31%20PM.png",
    "cg3": BASE + "media/graphics/architecture/ChatGPT%20Image%20Feb%2028%2C%202026%2C%2002_37_12%20PM.png",
    "cg4": BASE + "media/graphics/architecture/ChatGPT%20Image%20Feb%2028%2C%202026%2C%2002_38_41%20PM.png",
    "la1": BASE + "media/graphics/living-age/ChatGPT%20Image%20Mar%206%2C%202026%2C%2002_32_21%20PM.png",
    "la2": BASE + "media/graphics/living-age/ChatGPT%20Image%20Mar%206%2C%202026%2C%2002_34_47%20PM.png",
}

def fetch(url):
    try:
        req = urllib.request.Request(url, headers={"User-Agent":"Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=15) as r:
            return io.BytesIO(r.read())
    except Exception as e:
        print(f"  skip {url.split('/')[-1][:40]}: {e}"); return None

print("Downloading images...")
I = {k: fetch(u) for k,u in URLS.items()}

W, H = Inches(13.33), Inches(7.5)

def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    b = s.shapes.add_shape(1,0,0,W,H)
    b.fill.solid(); b.fill.fore_color.rgb = BG; b.line.fill.background()
    return s

def box(s,c,x,y,w,h):
    r=s.shapes.add_shape(1,x,y,w,h)
    r.fill.solid(); r.fill.fore_color.rgb=c; r.line.fill.background(); return r

def t(s, text, x,y,w,h, sz=16, bold=False, italic=False, color=WHT, align=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(x,y,w,h); tb.word_wrap=True
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=italic
    r.font.color.rgb=color; r.font.name="Calibri"

def ml(s, lines, x,y,w,h):
    tb=s.shapes.add_textbox(x,y,w,h); tb.word_wrap=True
    tf=tb.text_frame; tf.word_wrap=True
    for i,(txt,col,sz,b,it,sp) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_before=Pt(sp)
        r=p.add_run(); r.text=txt
        r.font.size=Pt(sz); r.font.bold=b; r.font.italic=it
        r.font.color.rgb=col; r.font.name="Calibri"

def bl(s, items, x,y,w,h):
    tb=s.shapes.add_textbox(x,y,w,h); tb.word_wrap=True
    tf=tb.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_before=Pt(6)
        if len(item)==2:
            r1=p.add_run(); r1.text=f"• {item[0]}:  "
            r1.font.bold=True; r1.font.color.rgb=GOLD
            r1.font.size=Pt(14); r1.font.name="Calibri"
            r2=p.add_run(); r2.text=item[1]
            r2.font.color.rgb=WHT; r2.font.size=Pt(13); r2.font.name="Calibri"
        else:
            r=p.add_run(); r.text=f"• {item[0]}"
            r.font.color.rgb=WHT; r.font.size=Pt(13); r.font.name="Calibri"

def pic(s,data,x,y,w,h=None):
    if data is None: return
    data.seek(0)
    s.shapes.add_picture(data,x,y,w,h) if h else s.shapes.add_picture(data,x,y,w)

def div(s,y=Inches(2.0)):
    r=s.shapes.add_shape(1,Inches(0.6),y,Inches(12.1),Pt(1))
    r.fill.solid(); r.fill.fore_color.rgb=DIM; r.line.fill.background()

def hdr(s,tag,title,sub=None):
    t(s,tag,Inches(0.6),Inches(0.35),Inches(10),Inches(0.4),sz=11,bold=True,color=TEAL)
    t(s,title,Inches(0.6),Inches(0.72),Inches(12),Inches(1.1),sz=34,bold=True,italic=True,color=GOLD)
    if sub: t(s,sub,Inches(0.6),Inches(1.55),Inches(10),Inches(0.45),sz=14,color=MUT)
    div(s)

def brandtag(s):
    box(s,GRN,Inches(0.6),Inches(0.3),Inches(2.6),Inches(0.35))
    t(s,"SPIRAL STEWARD",Inches(0.7),Inches(0.33),Inches(2.4),Inches(0.3),sz=11,bold=True,color=WHT)
    box(s,RGBColor(0x1c,0x40,0x4a),Inches(3.35),Inches(0.3),Inches(2.4),Inches(0.35))
    t(s,"LIVING WORKS",Inches(3.45),Inches(0.33),Inches(2.2),Inches(0.3),sz=11,bold=True,color=TEAL)

def note(s,text):
    s.notes_slide.notes_text_frame.text=text

prs=Presentation()
prs.slide_width=W; prs.slide_height=H

# ── SLIDE 1: TITLE ───────────────────────────────────────────────────────────
s=slide(prs)
box(s,RGBColor(0x0a,0x0e,0x14),0,0,Inches(7.0),H)
t(s,"SPIRAL STEWARD",Inches(0.7),Inches(1.1),Inches(6.2),Inches(0.9),sz=44,bold=True,italic=True,color=GOLD)
t(s,"×",Inches(0.7),Inches(1.95),Inches(0.8),Inches(0.55),sz=28,color=DIM)
t(s,"Living Works",Inches(1.4),Inches(1.92),Inches(5.5),Inches(0.6),sz=28,italic=True,color=TEAL)
r=s.shapes.add_shape(1,Inches(0.7),Inches(2.65),Inches(5.8),Pt(2))
r.fill.solid(); r.fill.fore_color.rgb=GRN; r.line.fill.background()
t(s,"The philosophy that living systems are\nthe right design paradigm — and the\ntechnology platform that proves it.",
  Inches(0.7),Inches(2.85),Inches(6.0),Inches(1.6),sz=19,color=WHT)
t(s,"Spiral Steward is the mission.\nLiving Works is the product.",
  Inches(0.7),Inches(4.55),Inches(6.0),Inches(0.8),sz=15,italic=True,color=MUT)
t(s,"sage-clokey.github.io  ·  Sage Clokey  ·  UCSC BME Bioinformatics",
  Inches(0.7),Inches(6.75),Inches(6.0),Inches(0.4),sz=11,color=DIM)
pic(s,I["spiral_house"],Inches(7.1),Inches(0.0),Inches(6.23),H)
note(s,"Two names because there are two things being built — and neither works without the other.\n\nSpiral Steward is the philosophy: living systems — organisms, spirals, mutualistic networks — are the superior design paradigm. Not as metaphor. As technical specification.\n\nLiving Works is the technology: the first Morphological BioCAD platform that lets you design with biological growth processes rather than geometric shapes.\n\nMost technology companies have a product but no worldview. Most philosophical movements have a worldview but no product. We have both. That is the moat.")

# ── SLIDE 2: PHILOSOPHY ──────────────────────────────────────────────────────
s=slide(prs)
brandtag(s)
t(s,"The Philosophy Behind the Product",Inches(0.6),Inches(0.82),Inches(12),Inches(0.55),sz=30,bold=True,italic=True,color=GOLD)
div(s,Inches(1.48))
ml(s,[
    ("There are two kinds of builders.",GOLD,18,True,False,0),
    ("The Mason carves the world into blocks and assembles them into structures. Every machine, every skyscraper, every drug molecule designed on a screen — imposed, specified, brittle at the edges.",MUT,13,False,False,10),
    ("The Steward reads the landscape, changes the conditions, and lets form emerge. A coral reef. A forest canopy. A developing embryo. No blueprint. No maintenance team. Half a billion years of successful iteration.",WHT,13,False,False,10),
    ("Evolution and markets are the same system.",GOLD,16,True,False,14),
    ("Both are decentralized, emergent, and capable of aggregating distributed information without central direction. Mutualism is voluntary trade — positive-sum, consent-based, fitness-increasing for both. The most stable systems in nature and economics are built on mutualism.",MUT,13,False,False,8),
    ('"The greatest integration in evolutionary history was not designed. Mitochondria were recruited through a relationship."',TEAL,14,True,True,10),
],Inches(0.6),Inches(1.6),Inches(6.8),Inches(5.6))
pic(s,I["block_vs_spiral"],Inches(7.7),Inches(1.6),Inches(5.3),Inches(4.2))
t(s,"The Mason imposes. The Steward cultivates.",Inches(7.7),Inches(5.88),Inches(5.3),Inches(0.4),sz=11,italic=True,color=DIM,align=PP_ALIGN.CENTER)
note(s,"Deliver this like a manifesto, not a market analysis.\n\nThree parts:\n1. The mason's model is wrong for complex systems — it optimizes single variables and fails at the edges.\n2. Living systems operate on fundamentally different principles — mutualism, emergence, self-organization — and produce superior outcomes.\n3. Evolution and markets are structurally identical emergent systems. This is the framework that tells us what interventions work.\n\nPunchline: we are building the first design platform grounded in the steward's model. Not because it is beautiful — because it is correct.")

# ── SLIDE 3: TEAM ────────────────────────────────────────────────────────────
s=slide(prs)
brandtag(s)
hdr(s,"03 · TEAM","Who Is Building This","Founder background, qualifications, and the story")
box(s,CARD,Inches(0.6),Inches(2.15),Inches(3.6),Inches(4.75))
t(s,"SAGE CLOKEY",Inches(0.78),Inches(2.3),Inches(3.2),Inches(0.5),sz=20,bold=True,color=GOLD)
t(s,"Founder",Inches(0.78),Inches(2.78),Inches(3.2),Inches(0.35),sz=13,color=TEAL)
t(s,"Senior · BME Bioinformatics\nUC Santa Cruz\n\nsage-clokey.github.io\ngithub.com/Sage-Clokey",Inches(0.78),Inches(3.18),Inches(3.1),Inches(2.6),sz=13,color=MUT)
bl(s,[
    ("The philosophy","Developed a complete framework connecting living systems, market theory, and design — published as the Spiral Steward manifesto. This is the product roadmap as first principles."),
    ("The science","Morphogenetic simulation, population genetics (Wright-Fisher), Sanger sequencing — BME bioinformatics at UCSC with direct research experience."),
    ("The engineering","18 deployed projects: Python, React, PyTorch, FastAPI, Kotlin (Android). Full-stack from mobile apps to ML growth simulations."),
    ("The track record","Built and shipped everything in this deck while carrying a full course load. Hungry. Not waiting for permission."),
],Inches(4.45),Inches(2.15),Inches(8.65),Inches(4.75))
note(s,"The story matters more than the credentials.\n\nBecame obsessed with one question: why do we design everything as a machine when the most successful systems on Earth are organisms? That question did not stay abstract — built the philosophy into a published manifesto, built morphogenesis simulations in PyTorch, deployed 18 projects publicly.\n\nThe unique qualification: a coherent worldview AND the ability to ship. Most builders lack the philosophy. Most philosophers cannot ship.")

# ── SLIDE 4: DEMO ────────────────────────────────────────────────────────────
s=slide(prs)
brandtag(s)
hdr(s,"04 · DEMO","Living Works — The Platform in Action","Natural language → growth algorithm → 3D morphogenetic output")
for i,(key,cap) in enumerate([("cg1","Fibonacci spiral house"),("cg2","Living facade — self-regulating building skin"),("mycelium_dome","Mycelium dome — fungal network as structural form"),("branching","Branching tree — vascular logic as load-bearing design")]):
    x=Inches(0.5)+(i%2)*Inches(6.35); y=Inches(2.15)+(i//2)*Inches(2.55)
    pic(s,I[key],x,y,Inches(6.1),Inches(2.3))
    t(s,cap,x,y+Inches(2.1),Inches(6.0),Inches(0.35),sz=10,color=DIM)
t(s,"Live: sage-clokey.github.io/Living-works-by-the-word/",Inches(0.5),Inches(7.12),Inches(12.5),Inches(0.3),sz=10,color=DIM)
note(s,"Key framing: in every other design tool, you draw the shape. In Living Works, you describe the growth conditions and the shape emerges — exactly as it does in biology.\n\nThe Spiral Steward philosophy is embedded in the product itself: you do not specify the form, you specify the conditions. The form is what the growth process finds.\n\nPoint to the live site for the technical audience.")

# ── SLIDE 5: MARKET ──────────────────────────────────────────────────────────
s=slide(prs)
brandtag(s)
hdr(s,"05 · MARKET","A $55B+ Uncontested Opportunity","Four converging industries — one platform connecting them all")
for i,(name,size,desc) in enumerate([
    ("Synthetic Biology","$20B+","Organism engineering for pharma, agriculture, materials"),
    ("AI Drug Design","$10B+","Generative protein and antibody modeling"),
    ("Bioinformatics Software","$15B+","Biological data management and analysis"),
    ("Living Materials","Emerging","Bio-grown buildings, mycelium composites, bacterial concrete"),
]):
    x=Inches(0.5)+i*Inches(3.1)
    box(s,CARD,x,Inches(2.15),Inches(2.95),Inches(3.3))
    t(s,size,x+Inches(0.18),Inches(2.28),Inches(2.6),Inches(0.72),sz=26,bold=True,color=GOLD)
    t(s,name,x+Inches(0.18),Inches(2.95),Inches(2.6),Inches(0.45),sz=13,bold=True,color=WHT)
    t(s,desc,x+Inches(0.18),Inches(3.4),Inches(2.6),Inches(1.0),sz=12,color=MUT)
box(s,ELEV,Inches(0.5),Inches(5.65),Inches(12.3),Inches(1.65))
t(s,"Why Spiral Steward creates a category, not just a product:",Inches(0.7),Inches(5.75),Inches(11),Inches(0.4),sz=14,bold=True,color=TEAL)
t(s,"Every market above is spending heavily on tools — and every existing tool treats biology as data to manage or molecules to assemble. None operate at the morphological scale: the growth process, the living form, the organism as designed output. Spiral Steward's philosophy defines that category. Living Works occupies it.",Inches(0.7),Inches(6.18),Inches(12.0),Inches(1.0),sz=13,color=MUT)
note(s,"Two layers:\n\n1. The numbers: $55B+ across four industries with no morphological design layer.\n\n2. The philosophy point: Spiral Steward is not entering an existing market — it is defining a new one. 'Morphological BioCAD' does not exist yet because no one had the philosophical framework to see it as a category. We do.\n\nThe living materials market is the fastest-growing and most uncontested. A firm designing a mycelium pavilion has no software today.")

# ── SLIDE 6: SOLUTION ────────────────────────────────────────────────────────
s=slide(prs)
brandtag(s)
hdr(s,"06 · SOLUTION","Philosophy Into Product","How Spiral Steward's principles become Living Works' four technology pillars")
ml(s,[
    ("The Spiral Steward principle:",TEAL,14,True,False,0),
    ("Do not specify the form. Specify the growth conditions.",WHT,17,True,True,4),
],Inches(0.6),Inches(2.15),Inches(5.8),Inches(0.9))
for i,(name,desc) in enumerate([
    ("Morphogenesis Engine","3D simulation of biological growth — branching systems, reaction-diffusion, spiral geometry, and biomechanical forces. The computational core of the steward's model."),
    ("Natural Language Bio Design","Type what you want to grow. The platform converts descriptions into growth algorithms — the mason specifies shape, the steward specifies conditions."),
    ("BioCAD Platform","Computer-aided design for organism blueprinting — the first CAD tool where the object builds itself rather than being drawn."),
    ("Living Materials Library","Curated database of biological building blocks — coral, mycelium, vascular trees — with growth parameters and simulation outputs."),
]):
    y=Inches(3.18)+i*Inches(1.0)
    box(s,CARD,Inches(0.6),y,Inches(5.8),Inches(0.92))
    t(s,f"0{i+1}  {name}",Inches(0.78),y+Inches(0.07),Inches(5.4),Inches(0.35),sz=14,bold=True,color=GOLD)
    t(s,desc,Inches(0.78),y+Inches(0.42),Inches(5.4),Inches(0.44),sz=11,color=MUT)
pic(s,I["mycelium_plan"],Inches(6.75),Inches(2.15),Inches(6.35),Inches(5.1))
note(s,"Bridge the philosophy to the product for each pillar:\n- Morphogenesis Engine: this is what 'specifying growth conditions' looks like computationally\n- NL Bio Design: the interface that makes the steward model accessible to non-engineers\n- BioCAD: what CAD looks like when biology is the design medium\n- Living Materials Library: not CAD blocks — living growth templates\n\nAutodesk gives you drawing tools. Unity gives you a physics engine. Living Works gives you a growth engine. That is a new primitive.")

# ── SLIDE 7: BUSINESS MODEL ──────────────────────────────────────────────────
s=slide(prs)
brandtag(s)
hdr(s,"07 · BUSINESS MODEL","Four Revenue Streams","SaaS subscriptions · AI design services · IP licensing · philosophy curriculum")
for i,(tier,price,clr,desc) in enumerate([
    ("Academic","$50–200/mo",TEAL,"Researchers and university labs — product-market fit layer"),
    ("Startup","$500/mo",GOLD,"Synthetic biology and living materials companies"),
    ("Enterprise","$2,000+/mo",WHT,"Pharma, architecture firms, government agencies"),
]):
    x=Inches(0.5)+i*Inches(3.1)
    box(s,CARD,x,Inches(2.15),Inches(2.95),Inches(2.15))
    t(s,tier,x+Inches(0.18),Inches(2.25),Inches(2.6),Inches(0.4),sz=16,bold=True,color=clr)
    t(s,price,x+Inches(0.18),Inches(2.62),Inches(2.6),Inches(0.48),sz=22,bold=True,color=WHT)
    t(s,desc,x+Inches(0.18),Inches(3.08),Inches(2.6),Inches(0.85),sz=11,color=MUT)
for i,(name,desc) in enumerate([
    ("AI Design Services","Custom organism and material design for pharma and architecture. High-margin consulting revenue."),
    ("IP Licensing","Patents on biological growth algorithms — long-term recurring revenue from pharma and materials partners."),
    ("Spiral Steward Curriculum","Philosophy and framework licensed to universities and design schools as a teaching standard. Brand revenue no competitor can replicate."),
]):
    x=Inches(0.5)+i*Inches(4.2)
    box(s,ELEV,x,Inches(4.5),Inches(3.95),Inches(2.8))
    t(s,name,x+Inches(0.18),Inches(4.62),Inches(3.6),Inches(0.4),sz=14,bold=True,color=GOLD)
    t(s,desc,x+Inches(0.18),Inches(5.06),Inches(3.6),Inches(2.0),sz=12,color=MUT)
note(s,"Four streams — make the fourth one land.\n\nSaaS is standard. Start with academics, expand to enterprise.\n\nAI design services: a company says 'design us a mycelium insulation material' and we deliver the growth specification.\n\nIP licensing: novel biological growth algorithms get patented and licensed.\n\nSpiral Steward curriculum licensing is the most defensible. When Spiral Steward becomes the standard philosophy for biological design — the way Bauhaus became the standard for industrial design — that is brand licensing revenue no competitor can replicate. The philosophy is the moat.")

# ── SLIDE 8: CUSTOMERS ───────────────────────────────────────────────────────
s=slide(prs)
brandtag(s)
hdr(s,"08 · CUSTOMERS","Who Uses Living Works","Technology buyers and philosophy adopters — two customer types, one platform")
pic(s,I["la2"],Inches(8.8),Inches(2.15),Inches(4.3),Inches(4.85))
bl(s,[
    ("Academic researchers","~15,000 university labs globally in BME, synthetic biology, biomaterials. Beachhead market. Already embedded at UCSC."),
    ("Biotech startups","Living materials, tissue engineering, sustainable construction. High willingness to pay."),
    ("Architecture firms","Exploring bio-grown structures and living facades. Large market, zero existing tooling."),
    ("Pharma & drug design","AI drug discovery companies needing morphological simulation at tissue and organ scale."),
    ("Philosophy adopters","Design schools, MBA programs, policy institutes adopting Spiral Steward as a teaching framework — curriculum licensing pipeline."),
    ("Acquisition path","Phase 1: UCSC lab network. Phase 2: SBIR credibility opens federal buyers. Phase 3: SynBioBeta, iGEM, SXSW Bio."),
],Inches(0.5),Inches(2.15),Inches(8.0),Inches(5.1))
note(s,"Two types of customer — make the distinction clear.\n\nTechnology customers use Living Works to design with biology — SaaS and design services.\n\nPhilosophy customers adopt the Spiral Steward framework to teach and research — curriculum licensing.\n\nThis second customer category is what makes this defensible. You cannot copy a philosophy. You can copy a product.\n\nThe beachhead is academic — already embedded in it. First customers are not hypothetical.")

# ── SLIDE 9: COMPETITION ─────────────────────────────────────────────────────
s=slide(prs)
brandtag(s)
hdr(s,"09 · COMPETITION","Why No One Else Is Here","The technical gap + the philosophical moat")
cols_x=[Inches(0.5),Inches(3.1),Inches(7.85)]
cols_w=[Inches(2.5),Inches(4.6),Inches(5.25)]
for j,(h,cx,cw) in enumerate(zip(["Competitor","What They Do","What They Miss"],cols_x,cols_w)):
    box(s,GOLD,cx,Inches(2.15),cw-Inches(0.08),Inches(0.42))
    t(s,h,cx+Inches(0.12),Inches(2.2),cw-Inches(0.2),Inches(0.35),sz=13,bold=True,color=BG)
for i,(comp,does,miss) in enumerate([
    ("Benchling","Lab data management, LIMS, sequence editing","No simulation. No growth. Manages biology, doesn't design it."),
    ("Ginkgo Bioworks","Organism engineering at genetic level","Designs genes, not growth. No morphological output."),
    ("Autodesk BioCAD","Early molecular design tools","Protein/molecular scale only — no organism-level morphogenesis."),
    ("Rhino + Grasshopper","Parametric architecture design","Geometric, not biological. You specify shape — not growth."),
    ("Codon Devices","DNA synthesis, metabolic pathway design","Genetic level only. No morphological output."),
]):
    ry=Inches(2.6)+i*Inches(0.77)
    bg=CARD if i%2==0 else BG
    for cx,cw in zip(cols_x,cols_w):
        box(s,bg,cx,ry,cw-Inches(0.08),Inches(0.73))
    t(s,comp,cols_x[0]+Inches(0.12),ry+Inches(0.08),cols_w[0]-Inches(0.2),Inches(0.58),sz=13,bold=True,color=TEAL)
    t(s,does,cols_x[1]+Inches(0.12),ry+Inches(0.05),cols_w[1]-Inches(0.2),Inches(0.66),sz=11,color=MUT)
    t(s,miss,cols_x[2]+Inches(0.12),ry+Inches(0.05),cols_w[2]-Inches(0.2),Inches(0.66),sz=11,color=WHT)
box(s,ELEV,Inches(0.5),Inches(6.52),Inches(12.3),Inches(0.75))
t(s,"Our moat: Living Works is the only platform at the morphological scale. Spiral Steward is the only coherent philosophy for why that scale matters. A competitor can build a growth engine. They cannot build 10 years of philosophical development.",Inches(0.68),Inches(6.6),Inches(12.0),Inches(0.62),sz=13,bold=True,color=GOLD)
note(s,"Address Benchling proactively — it will come up. Benchling manages lab data. We generate design outputs. Completely different product.\n\nTechnical moat: we operate at the morphological scale — unoccupied by every competitor.\n\nPhilosophical moat: Spiral Steward defines why the morphological scale matters. A competitor can replicate the technology. They cannot replicate the philosophical development that is already published and being taught.")

# ── SLIDE 10: FINANCIALS ─────────────────────────────────────────────────────
s=slide(prs)
brandtag(s)
hdr(s,"10 · FINANCIAL OVERVIEW","Revenue Model & Projections","Conservative estimates based on comparable SaaS biology platforms")
for i,(phase,rev,clr,desc) in enumerate([
    ("Year 1–2 · Foundation","$0–50K ARR",TEAL,"10–50 academic users. SBIR Phase I ($275K non-dilutive). Spiral Steward curriculum pilots."),
    ("Year 3–5 · Expansion","$500K–2M ARR",GOLD,"100+ startups, 5–10 enterprise accounts. AI design services. SBIR Phase II ($1.5M). Curriculum licensing."),
    ("Year 5–10 · Scale","$10M+ ARR",WHT,"Platform scale. IP licensing revenue. Spiral Steward in 50+ university curricula."),
]):
    x=Inches(0.5)+i*Inches(4.2)
    box(s,CARD,x,Inches(2.15),Inches(3.95),Inches(2.95))
    t(s,phase,x+Inches(0.18),Inches(2.25),Inches(3.6),Inches(0.62),sz=13,bold=True,color=clr,align=PP_ALIGN.LEFT)
    t(s,rev,x+Inches(0.18),Inches(2.83),Inches(3.6),Inches(0.55),sz=22,bold=True,color=WHT)
    t(s,desc,x+Inches(0.18),Inches(3.36),Inches(3.6),Inches(1.55),sz=11,color=MUT)
bl(s,[
    ("Burn rate","~$30–50K/month with 3–4 person founding team"),
    ("Runway","$500K seed = 12–18 months to Series A milestones"),
    ("Non-dilutive","SBIR Phase I ($275K) + Phase II ($1.5M) — significant runway without equity dilution"),
    ("Curriculum revenue","Defensible revenue stream independent of platform scale"),
],Inches(0.5),Inches(5.35),Inches(12.5),Inches(1.95))
note(s,"Pre-revenue, pre-incorporation. Projections grounded in comparable SaaS biology platforms.\n\nKey financial story: the SBIR pathway — $275K Phase I non-dilutive, $1.5M Phase II. That is $1.775M without touching equity.\n\nSpiral Steward curriculum licensing is revenue most biotech companies do not have. Once in curricula, it becomes self-reinforcing.")

# ── SLIDE 11: FUNDING ────────────────────────────────────────────────────────
s=slide(prs)
brandtag(s)
hdr(s,"11 · FUNDING","The Ask","Seed round to build the foundation and validate product-market fit")
pic(s,I["sunflower"],Inches(9.1),Inches(2.15),Inches(4.0),Inches(4.85))
t(s,"$500,000",Inches(0.65),Inches(2.2),Inches(6),Inches(1.0),sz=52,bold=True,color=GOLD)
t(s,"Seed Round",Inches(0.65),Inches(3.1),Inches(5),Inches(0.4),sz=18,color=TEAL)
bl(s,[
    ("Platform development","40% — Morphogenesis engine, BioCAD interface, NL design layer"),
    ("Team","35% — Lead engineer + biological scientist co-founder"),
    ("Go-to-market","15% — Academic outreach, conferences, SBIR application support"),
    ("Infrastructure & legal","10% — Cloud compute, IP filings, Delaware C-Corp formation"),
],Inches(0.65),Inches(3.68),Inches(8.1),Inches(2.9))
box(s,ELEV,Inches(0.65),Inches(6.72),Inches(8.1),Inches(0.58))
t(s,"Raised to date: $0  ·  Non-dilutive pipeline: SBIR Phase I ($275K) in progress",Inches(0.82),Inches(6.8),Inches(7.8),Inches(0.4),sz=12,color=TEAL)
note(s,"Direct and honest: this is the first check. Pre-revenue, pre-incorporation.\n\nSBIR pathway is the key story — $275K Phase I non-dilutive, $1.5M Phase II. Federal credibility accelerates enterprise sales.\n\n$500K for 12–18 months to hit four Series A milestones: working morphogenesis engine, 50+ paying academic users, first enterprise pilot, SBIR Phase I award.")

# ── SLIDE 12: MILESTONES + LEGAL ─────────────────────────────────────────────
s=slide(prs)
brandtag(s)
hdr(s,"12 · MILESTONES & LEGAL","The Road Ahead","Technology milestones, philosophy milestones, and legal structure — all three tracked")
for i,(phase,years,name,clr,mstones) in enumerate([
    ("Phase 1","Years 1–2","Foundation",TEAL,[
        "Morphogenesis engine v1 — branching, reaction-diffusion, spiral",
        "BioCAD interface — NL prompt to 3D growth output",
        "50+ paying academic users",
        "SBIR Phase I award ($275K non-dilutive)",
        "UCSC lab partnership",
        "Spiral Steward curriculum in 1+ university",
        "Delaware C-Corp formation",
    ]),
    ("Phase 2","Years 3–5","Expansion",GOLD,[
        "Living Materials module — mycelium, bio-concrete",
        "Enterprise pilot with pharma or architecture firm",
        "SBIR Phase II award ($1.5M non-dilutive)",
        "Series A round ($3–5M)",
        "100+ startup users, 5+ enterprise contracts",
        "Spiral Steward in 10+ university curricula",
        "First peer-reviewed publication using Living Works",
    ]),
    ("Phase 3","Years 5–10","Platform Standard",WHT,[
        "Programmable ecosystems — multi-species growth",
        "Living infrastructure standardization",
        "IP licensing from pharma and materials partners",
        "Spiral Steward = industry philosophy standard",
        "Living Works = Autodesk of morphological biology",
    ]),
]):
    x=Inches(0.4)+i*Inches(4.3)
    box(s,clr,x,Inches(2.15),Inches(4.05),Inches(0.56))
    t(s,f"{phase} · {name}",x+Inches(0.15),Inches(2.19),Inches(3.8),Inches(0.4),sz=13,bold=True,color=BG)
    t(s,years,x+Inches(0.15),Inches(2.68),Inches(3.8),Inches(0.28),sz=11,color=MUT)
    box(s,CARD,x,Inches(2.96),Inches(4.05),Inches(4.34))
    for j,m in enumerate(mstones):
        t(s,f"→  {m}",x+Inches(0.15),Inches(3.02)+j*Inches(0.58),Inches(3.78),Inches(0.53),sz=11,color=WHT if j<5 else MUT)
note(s,"Every phase has BOTH technology AND philosophy milestones — that is intentional.\n\nTechnology milestones prove the product works. Philosophy milestones (curriculum adoption, publications) prove the category is being defined and owned.\n\nLegal: pre-incorporation. Delaware C-Corp happens at first check or SBIR award. All IP owned by founder, transfers cleanly to company. Engaging UCSC Office of Research to confirm personal-time IP is not subject to university claims — standard for UC student-founded companies.\n\nLocation: UC Santa Cruz, California. Bay Area proximity for biotech and VC access.")

prs.save(OUT)
print(f"\nSaved: {OUT}")
print(f"Slides: {len(prs.slides)}")
